"""
Generate a 15-slide final-year project presentation (PPTX) following
the provided "Presentation guide.pdf" structure.

Output:
  Route_Optimization_Presentation.pptx
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_TITLE = "Field Sales Route Optimization System"
SUBTITLE = "On-demand agent routing using road graphs, priority scoring, and time-window constraints"


def _set_title(slide, title: str, subtitle: str | None = None) -> None:
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True

    if subtitle is not None:
        sub = slide.placeholders[1]
        sub.text = subtitle
        p = sub.text_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(80, 80, 80)


def _add_bullets(slide, title: str, bullets: list[str]) -> None:
    # Caller provides an already-created slide (Title and Content).
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(20)


def _add_table_lit_review(slide) -> None:
    slide.shapes.title.text = "Literature Review & Theoretical Framework"

    rows, cols = 4, 4
    left = Inches(0.8)
    top = Inches(1.7)
    width = Inches(8.6)
    height = Inches(2.2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    headers = ["Author", "Method", "Findings", "Gap"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)

    examples = [
        ("Study A", "Shortest-path routing", "Accurate road-based distances", "No per-agent constraints"),
        ("Study B", "VRP heuristics", "Fast route construction", "No shop time-windows"),
        ("This work", "Dijkstra + constrained NN", "Agent-ready routes + map UI", "—"),
    ]
    for i, row in enumerate(examples, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)

    # Gap statement
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), Inches(8.6), Inches(1.1)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(235, 245, 255)
    box.line.color.rgb = RGBColor(37, 99, 235)
    tf = box.text_frame
    tf.text = "Identified gap: routing that is both road-accurate and operationally feasible per agent (hours, time windows, home endpoints)."
    tf.paragraphs[0].font.size = Pt(14)


def _add_equation_box(slide, title: str, equation: str, note: str) -> None:
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = note
    p.font.size = Pt(18)

    eq = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.0), Inches(7.6), Inches(1.2)
    )
    eq.fill.solid()
    eq.fill.fore_color.rgb = RGBColor(245, 245, 245)
    eq.line.color.rgb = RGBColor(150, 150, 150)
    tf = eq.text_frame
    tf.clear()
    p2 = tf.paragraphs[0]
    p2.text = equation
    p2.font.size = Pt(26)
    p2.alignment = PP_ALIGN.CENTER


def _add_architecture(slide) -> None:
    slide.shapes.title.text = "Proposed Methodology & System Architecture"

    # Simple block diagram (left-to-right)
    blocks = [
        ("Input Excel\n(locations, HQ, agents)", RGBColor(226, 232, 240)),
        ("Priority scoring\n+ barrier (top-k)", RGBColor(219, 234, 254)),
        ("Road graph cache\n(OSMnx GraphML)", RGBColor(220, 252, 231)),
        ("Dijkstra matrices\n(length, travel_time)", RGBColor(254, 249, 195)),
        ("Constrained routing\n(time windows, hours)", RGBColor(254, 226, 226)),
        ("Web UI + Map\n(Folium/Leaflet)", RGBColor(237, 233, 254)),
    ]

    left = Inches(0.5)
    top = Inches(2.0)
    w = Inches(1.55)
    h = Inches(1.0)
    gap = Inches(0.25)

    shapes = []
    for i, (label, color) in enumerate(blocks):
        x = left + i * (w + gap)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.color.rgb = RGBColor(80, 80, 80)
        tf = rect.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        shapes.append(rect)

    # Arrows
    for i in range(len(shapes) - 1):
        x1 = shapes[i].left + shapes[i].width
        x2 = shapes[i + 1].left
        y = shapes[i].top + shapes[i].height // 2
        line = slide.shapes.add_connector(1, x1, y, x2, y)  # 1 = straight
        line.line.color.rgb = RGBColor(37, 99, 235)
        line.line.width = Pt(2)


def build_pptx(output_path: Path) -> None:
    prs = Presentation()
    prs.core_properties.title = PROJECT_TITLE

    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    _set_title(
        s1,
        PROJECT_TITLE,
        f"{SUBTITLE}\nStudent: ____________________   Reg No: ____________________\nSupervisor: ____________________   Date: {date.today().isoformat()}",
    )

    # Slide 2: Background
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s2,
        "Introduction & Research Background",
        [
            "Field sales teams must visit many shops daily under time and operational constraints",
            "Manual planning is inefficient and often ignores road reality, opening hours, and end-of-day home travel",
            "Goal: improve productivity, coverage, and expected revenue with automated, agent-ready routes",
        ],
    )

    # Slide 3: Problem & objectives
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s3,
        "Problem Statement & Research Objectives",
        [
            "Problem: generate feasible daily routes per agent that follow roads and respect hours + shop time windows",
            "RO1: compute road-accurate travel distance/time using Dijkstra on OpenStreetMap road graphs",
            "RO2: prioritize shops using multi-factor sales signals and select top-k per agent (priority barrier)",
            "RO3: produce a usable artifact (web UI) that gives each agent a clear route map and ordered stops",
        ],
    )

    # Slide 4: Literature table
    s4 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    _add_table_lit_review(s4)

    # Slide 5: Data acquisition & preprocessing
    s5 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s5,
        "Data Acquisition & Preprocessing Pipeline",
        [
            "Sources: locations.xlsx (shops + sales signals), hq.xlsx (start), agents.xlsx (home endpoints)",
            "Preprocessing: validate coordinates, handle missing optional fields with defaults, remove coordinate outliers (IQR)",
            "Feature preparation: normalize conversion rate and order value; compute recency score; encode shop size",
        ],
    )
    _add_equation_box(
        s5,
        "Data Acquisition & Preprocessing Pipeline",
        "z = (x − μ) / σ",
        "Standardization is used to put numeric features on comparable scales when needed.",
    )

    # Slide 6: EDA & feature engineering (text + placeholders)
    s6 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s6,
        "EDA & Feature Engineering (high-yield insights)",
        [
            "Investigated shop distribution per agent/sub-cluster and missingness in optional columns",
            "Checked correlations between conversion_rate, avg_order_value, and last_visit_days for scoring stability",
            "Engineered: recency = 1 − exp(−days / τ) and size score mapping (small/medium/large)",
        ],
    )

    # Slide 7: Architecture
    s7 = prs.slides.add_slide(prs.slide_layouts[5])
    _add_architecture(s7)

    # Slide 8: Model selection & training (here: algorithm selection)
    s8 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s8,
        "Algorithm Selection & Route Construction",
        [
            "Baseline: straight-line distances / naive ordering (not road-accurate, often infeasible)",
            "Selected: Dijkstra shortest paths on road graphs for distance/time matrices",
            "Route ordering: constrained nearest-neighbor heuristic for fast per-agent routing",
            "Setup: per agent, select top-k shops then build route HQ → shops → Home",
        ],
    )

    # Slide 9: Hyperparameters / loss functions (here: tunables + constraints)
    s9 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s9,
        "Optimization Setup (Tuning & Constraints)",
        [
            "Tunable parameters: daily_shop_target (k), max_work_hours, work_start_hour, priority weights",
            "Hard constraints: shop time windows, service time per shop, must reach home before work end",
            "Graph caching: GraphML cache reused across runs for faster on-demand generation",
        ],
    )
    _add_equation_box(
        s9,
        "Optimization Setup (Tuning & Constraints)",
        "Priority = 0.35·CR + 0.30·AOV + 0.20·Recency + 0.15·Size",
        "Multi-factor priority scoring selects the most valuable shops before routing.",
    )

    # Slide 10: Quantitative results & metrics (use known dataset stats + placeholders)
    s10 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s10,
        "Quantitative Results & Performance Metrics",
        [
            "Dataset: 1,575 shops, 15 agents, 4 sub-clusters (from system logs)",
            "Per-agent outputs: shops selected (top-k), shops visited, shops dropped (time/barrier), distance (km), travel time (min)",
            "Example (Chris Kibet): 26 shops visited, 110.16 km, 131.3 min travel, est. revenue KSh 16,941.48",
        ],
    )

    # Slide 11: Discussion & error analysis
    s11 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s11,
        "Discussion & Error Analysis",
        [
            "Dropped stops occur mainly due to time constraints (cannot finish before end-of-day) or priority barrier",
            "Failure modes: unreachable shop coordinates (no path), overly tight shop hours, unrealistic service times",
            "Guardrails: coordinate validation, outlier removal, and fallback handling (straight line if no segment path)",
        ],
    )

    # Slide 12: Implementation & tech stack
    s12 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s12,
        "System Implementation & Tech Stack",
        [
            "Backend: Python (pandas, networkx, osmnx), Dijkstra shortest paths, constrained nearest-neighbor",
            "Mapping: Folium/Leaflet interactive HTML, numbered stops, direction arrows",
            "Web app: Flask API + UI (agent selects name, route generated on demand ~30–60s with cached graphs)",
            "Outputs: per-agent map HTML + Excel KPI summary report",
        ],
    )

    # Slide 13: Live demo placeholder + contingency
    s13 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s13,
        "Live Demonstration / Artifact Walkthrough",
        [
            "Demo scenario: agent selects name → route generates → follow numbered stops on the map",
            "Show: step-by-step directions list + map arrows + stop popups (service time & next stop)",
            "Contingency: embed a 45-second backup screen recording in this slide if needed",
        ],
    )

    # Slide 14: Contributions/limitations/future work
    s14 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s14,
        "Contributions, Limitations, and Future Work",
        [
            "Contributions: road-accurate routing, multi-factor shop prioritization, operational constraints, agent-facing web UI",
            "Limitations: heuristic route may be 15–25% longer than optimum; no traffic modeling; weights are fixed (no learning)",
            "Future work: add traffic/time-dependent speeds, improved solver (OR-Tools), feedback loop to learn scoring weights",
        ],
    )

    # Slide 15: References & acknowledgments
    s15 = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        s15,
        "References & Acknowledgments",
        [
            "Boeing, G. (2017). OSMnx: New methods for acquiring, constructing, analyzing, and visualizing complex street networks. Computers, Environment and Urban Systems.",
            "Dijkstra, E. W. (1959). A note on two problems in connexion with graphs. Numerische Mathematik.",
            "Holland, J. H. (1992). Adaptation in Natural and Artificial Systems (heuristics background).",
            "Thank you to the supervisor and panel for guidance and feedback.",
        ],
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    out = Path(__file__).resolve().parent / "Route_Optimization_Presentation.pptx"
    build_pptx(out)
    print(f"Saved: {out}")

